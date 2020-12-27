import os
import pptx
from pptx.util import Inches
import argparse
import os
import cv2
import sys
from PyQt5.QtWidgets import QApplication
from PIL import Image

def get_screen_dpi() -> float:
    '''Gets the DPI of a user.

    Returns:
        float -> The DPI of one's monitor
    '''
    # https://stackoverflow.com/questions/54271887/calculate-screen-dpi
    app = QApplication(sys.argv)
    screen = app.screens()[0]
    dpi = screen.physicalDotsPerInch()
    app.quit()

    return dpi

def pixels_to_inches(pix: int, dpi: float) -> int:
    '''Converts pixels to inches

    Arguments:
        pix (int) -> The pixel count (the number being converted)
        dpi (float) -> The DPI of one's monitor

    Returns:
        int -> The converted value
    '''
    return pix / dpi

def split_frames(filepath: str, verbose=True) -> int:
    '''Splits a given video into indivitual frames and saves it in the frames directory

    Arguments:
        filepath (str): The filepath of the video
        verbose (bool): Should there be verbose when FFMPEG splits the frames

    Returns:
        int -> 1 if it was successful. 0 if it failed
    '''
    if verbose:
        os.system(f"ffmpeg -i {filepath} -r 24/1 frames/%d.jpg")
    else:
        os.system(f"ffmpeg -hide_banner -loglevel panic -i {filepath} -r 24/1 frames/%d.jpg")
    return 1

def get_image_size(filepath: str) -> tuple:
    '''Get's the image dimentions as inches

    Arguments:
        filepath (str): The filepath of the image

    Returns:
        tuple: A tuple formatted as (width, height)
    '''
    im = cv2.imread(filepath)
    h, w, d = im.shape
    dpi = get_screen_dpi()

    h_inches = pixels_to_inches(h, dpi)
    w_inches = pixels_to_inches(w, dpi)
    return h_inches, w_inches

def make_slideshow(output_file_name="output.pptx") -> int:
    '''Makes the slideshow. Assumes that there are images in the frames directory

    Returns:
        int: Returns a 1 on success
    '''
    image_height, image_width = get_image_size("frames/1.jpg")

    presentation = pptx.Presentation()
    presentation.slide_width = Inches(image_width)
    presentation.slide_height = Inches(image_height)

    blank_slide = presentation.slide_layouts[6]

    for i in range(1, len(os.listdir("frames")) + 1, 1):
        image_path = f"frames/{i}.jpg"
        slide = presentation.slides.add_slide(blank_slide)
        slideshow_picture = slide.shapes.add_picture(image_path, width=Inches(image_width), height=Inches(image_height), left=0, top=0)

    presentation.save(output_file_name)

    return 1

def clear_frames_dir() -> int:
    '''Clears all the images in the frames directory

    Returns:
        int -> 1 if it was successful. 0 if it failed
    '''
    os.chdir("frames")
    for file in os.listdir(): os.remove(file)
    return 1


def main():
    args_obj = argparse.ArgumentParser()
    args_obj.add_argument("--filename", type=str, required=True, help="The filename that you would like to convert to a slideshow.")
    args_obj.add_argument("--ffmpeg_verbose", required=False, action='store_false', help="Would you like verbose in your FFMPEG program.")
    args_obj.add_argument("--keep_frames", required=False, action='store_true', help="Would you like to keep the split frames?")
    args_obj.add_argument("--output", required=False, default='output.pptx', help="The output name of your file.")
    args = args_obj.parse_args()

    filename = args.filename
    ffmpeg_verbose = args.ffmpeg_verbose
    keep_frames = args.keep_frames
    output = args.output

    split_frames(filename, verbose=ffmpeg_verbose)
    make_slideshow(output_file_name=output)

    if not keep_frames:
        clear_frames_dir()



if __name__ == "__main__":
    main()